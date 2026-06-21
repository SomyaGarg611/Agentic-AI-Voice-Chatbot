"""
Multi-format document extraction.

Turns raw bytes of many document types into a list of structured *blocks*:

  {"type": "prose", "text": str, "meta": {...}}
  {"type": "table", "header": [str], "rows": [[str]], "meta": {...}}

Prose blocks are later chunked by the recursive token-aware splitter.
Table blocks are chunked row-wise with the header repeated in every chunk,
so column context is never lost during retrieval.

Supported: PDF, DOCX, XLSX/XLSM, CSV, TSV, PPTX, TXT, MD, HTML, and a
UTF-8 fallback for anything else.
"""
import csv
import io
import logging
from typing import List, Dict, Any

logger = logging.getLogger(__name__)

Block = Dict[str, Any]


def extract_blocks(data: bytes, filename: str) -> List[Block]:
    name = (filename or "").lower()
    try:
        if name.endswith(".pdf"):
            return _pdf(data)
        if name.endswith((".docx",)):
            return _docx(data, filename)
        if name.endswith((".xlsx", ".xlsm")):
            return _xlsx(data, filename)
        if name.endswith(".csv"):
            return _delimited(data, filename, ",")
        if name.endswith(".tsv"):
            return _delimited(data, filename, "\t")
        if name.endswith(".pptx"):
            return _pptx(data, filename)
        if name.endswith((".html", ".htm")):
            return _html(data, filename)
        # txt, md, json, code, unknown → plain text
        return _plain(data, filename)
    except Exception as exc:
        logger.warning("Extraction failed for %s (%s) — falling back to plain text", filename, exc)
        return _plain(data, filename)


# ── PDF ─────────────────────────────────────────────────────────────
def _pdf(data: bytes) -> List[Block]:
    import fitz  # PyMuPDF
    blocks: List[Block] = []
    doc = fitz.open(stream=data, filetype="pdf")
    try:
        for page_no, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                blocks.append({"type": "prose", "text": text, "meta": {"page": page_no}})
    finally:
        doc.close()
    if not blocks:
        # Scanned/imaged PDF with no text layer
        logger.info("PDF has no extractable text layer")
    return blocks


# ── DOCX ────────────────────────────────────────────────────────────
def _docx(data: bytes, filename: str) -> List[Block]:
    import docx
    document = docx.Document(io.BytesIO(data))
    blocks: List[Block] = []

    # Paragraphs → one prose block (splitter handles sizing)
    paras = [p.text.strip() for p in document.paragraphs if p.text and p.text.strip()]
    if paras:
        blocks.append({"type": "prose", "text": "\n\n".join(paras), "meta": {}})

    # Tables → structured table blocks
    for t_idx, table in enumerate(document.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        rows = [r for r in rows if any(c for c in r)]
        if len(rows) >= 1:
            header, body = rows[0], rows[1:]
            blocks.append({"type": "table", "header": header, "rows": body,
                           "meta": {"table": t_idx + 1}})
    return blocks


# ── XLSX / XLSM ─────────────────────────────────────────────────────
def _xlsx(data: bytes, filename: str) -> List[Block]:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: List[Block] = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                rows.append(cells)
        if not rows:
            continue
        header, body = rows[0], rows[1:]
        blocks.append({"type": "table", "header": header, "rows": body,
                       "meta": {"sheet": ws.title}})
    wb.close()
    return blocks


# ── CSV / TSV ───────────────────────────────────────────────────────
def _delimited(data: bytes, filename: str, delim: str) -> List[Block]:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text), delimiter=delim)
    rows = [[c.strip() for c in r] for r in reader if any(cell.strip() for cell in r)]
    if not rows:
        return []
    header, body = rows[0], rows[1:]
    return [{"type": "table", "header": header, "rows": body, "meta": {}}]


# ── PPTX ────────────────────────────────────────────────────────────
def _pptx(data: bytes, filename: str) -> List[Block]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    blocks: List[Block] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            blocks.append({"type": "prose", "text": "\n".join(parts),
                           "meta": {"slide": slide_no}})
    return blocks


# ── HTML ────────────────────────────────────────────────────────────
def _html(data: bytes, filename: str) -> List[Block]:
    html = data.decode("utf-8", errors="replace")
    try:
        import trafilatura
        text = trafilatura.extract(html, include_comments=False, include_tables=True)
        if text and text.strip():
            return [{"type": "prose", "text": text.strip(), "meta": {}}]
    except Exception:
        pass
    return _plain(data, filename)


# ── Plain text fallback ─────────────────────────────────────────────
def _plain(data: bytes, filename: str) -> List[Block]:
    text = data.decode("utf-8", errors="replace").strip()
    return [{"type": "prose", "text": text, "meta": {}}] if text else []
